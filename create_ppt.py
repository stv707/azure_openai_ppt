import openai
import pptx
from pptx.util import Inches
from dotenv import load_dotenv
import os
import subprocess

# Load environment variables from .env file
load_dotenv()

# Define the path to the PowerPoint template
TEMPLATE_PATH = 'mytemp.pptx'

# Function to generate content using OpenAI API
def generate_content(prompt):
    openai.api_key = os.getenv('OPENAI_API_KEY')
    print(f"Sending prompt to OpenAI API: {prompt}")  # Debug information
    response = openai.ChatCompletion.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": "You are an assistant that helps create PowerPoint presentations."},
            {"role": "user", "content": prompt}
        ]
    )
    content = response['choices'][0]['message']['content']
    print(f"Received content from OpenAI API: {content}")  # Debug information
    return content

# Function to parse the content into title and bullet points / NOT performing, need to Parse More!! 

def parse_content(content):
    lines = content.split('\n')
    title = lines[0].strip('"')  # Remove double quotes from the title
    bullets = [line.strip('"') for line in lines[1:] if line.strip().startswith('-')][:4]  # Limit to 4 bullet points and remove double quotes
    return title, bullets


# Function to create PowerPoint slides using a template
def create_ppt(slide_data, template_path):
    prs = pptx.Presentation(template_path)

    for slide_info in slide_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Using slide layout 1 for title and content
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_info['title']
        content.text = '\n'.join(slide_info['bullets'])

    prs.save('presentation.pptx')

def main():
    # Collect information from the user
    presentation_title = input("Enter the title of the presentation: ")
    main_topic = input("Enter the main topic of the presentation: ")
    num_slides = int(input("Enter the number of slides: "))

    slide_data = []

    for i in range(num_slides):
        # Generate a subtopic for each slide
        subtopic_prompt = f"Generate a subtopic for slide {i + 1} of a PowerPoint presentation on the topic '{main_topic}'."
        subtopic = generate_content(subtopic_prompt).strip()

        # Generate content for each slide based on the subtopic
        content_prompt = f"Generate a title and 4 bullet points for a PowerPoint slide on the subtopic '{subtopic}'. Format it as follows: [Title of Slide]\n- [Bullet point 1]\n- [Bullet point 2]\n- [Bullet point 3]\n- [Bullet point 4]"
        content = generate_content(content_prompt)

        title, bullets = parse_content(content)

        slide_data.append({
            'title': title,
            'bullets': bullets
        })

    # Create the PowerPoint presentation using the template
    create_ppt(slide_data, TEMPLATE_PATH)

    print("Presentation created successfully as 'presentation.pptx'")

    # Ask if the user wants to open the presentation
    open_presentation = input("Shall we open the presentation? (Yes/No): ")
    if open_presentation.lower() in ['yes', 'y']:
        # Open the presentation
        if os.name == 'nt':  # For Windows
            os.startfile('presentation.pptx')
        else:  # For macOS/Linux
            subprocess.call(['open', 'presentation.pptx'])

if __name__ == "__main__":
    main()